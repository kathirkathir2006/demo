from pathlib import Path
from textwrap import dedent

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PPTX = Path("Python_Intro_Day1_YouTube.pptx")
OUT_NOTES = Path("Python_Intro_Day1_Speaker_Notes.md")


PALETTE = {
    "navy": RGBColor(10, 16, 34),
    "ink": RGBColor(20, 24, 39),
    "muted": RGBColor(100, 116, 139),
    "paper": RGBColor(248, 250, 252),
    "white": RGBColor(255, 255, 255),
    "cyan": RGBColor(34, 211, 238),
    "blue": RGBColor(59, 130, 246),
    "violet": RGBColor(139, 92, 246),
    "pink": RGBColor(236, 72, 153),
    "yellow": RGBColor(250, 204, 21),
    "green": RGBColor(34, 197, 94),
    "orange": RGBColor(251, 146, 60),
    "red": RGBColor(248, 113, 113),
    "slate": RGBColor(30, 41, 59),
    "code": RGBColor(15, 23, 42),
}


SLIDES = [
    {
        "kind": "cover",
        "section": "Cold open",
        "title": "Python Day 1",
        "subtitle": "Make the computer listen.",
        "kicker": "A cinematic intro for absolute beginners",
        "script": "Open with energy: Today is not about memorising syntax. Today we learn how to make a machine follow our words. Tell viewers they will write a tiny YouTube idea generator by the end.",
    },
    {
        "kind": "story",
        "section": "Cold open",
        "title": "Before Python, you already think like a programmer",
        "subtitle": "Daily life is full of tiny algorithms.",
        "cards": [
            ("Ordering tea", "If sugar is low, add more. If guests come, make extra."),
            ("Choosing a route", "Compare time, traffic, distance, then decide."),
            ("Editing a video", "Cut clip, add title, export, upload."),
        ],
        "script": "Normalize programming. A beginner is not starting from zero; they already make rules, decisions, and sequences every day.",
    },
    {
        "kind": "map",
        "section": "Roadmap",
        "title": "Today's episode has 3 acts",
        "subtitle": "Talk. Store. Ask.",
        "steps": [
            ("1", "Talk", "print() makes Python speak."),
            ("2", "Store", "Variables hold useful values."),
            ("3", "Ask", "input() makes the program interactive."),
        ],
        "script": "Set expectations. The whole day is about three superpowers: making output, saving information, and accepting user input.",
    },
    {
        "kind": "concept",
        "section": "Mindset",
        "title": "Programming is not magic",
        "subtitle": "It is a deal between you and a very fast, very literal intern.",
        "points": [
            "You give exact instructions.",
            "The computer follows them in order.",
            "If words are unclear, it complains with an error.",
        ],
        "script": "Use the 'literal intern' metaphor. Computers are powerful but not intelligent in the human sense; they need precise instructions.",
    },
    {
        "kind": "concept",
        "section": "What is Python?",
        "title": "Python is the friendly language layer",
        "subtitle": "It lets human ideas travel into machine actions.",
        "points": [
            "Readable enough for beginners.",
            "Powerful enough for automation, AI, web apps, and data.",
            "Used by creators to save time and build tools.",
        ],
        "script": "Avoid a boring history lecture. Position Python as a bridge from ideas to actions, especially for creator workflows.",
    },
    {
        "kind": "pipeline",
        "section": "How it runs",
        "title": "What happens when you run Python?",
        "subtitle": "Your code goes through a simple pipeline.",
        "steps": [
            ("You", "Write code"),
            ("Python", "Interprets line by line"),
            ("Computer", "Shows result"),
        ],
        "script": "Explain interpreter simply. Python reads code line by line and asks the computer to do things.",
    },
    {
        "kind": "setup",
        "section": "Setup",
        "title": "Day 1 setup: keep it simple",
        "subtitle": "You only need three places today.",
        "items": [
            ("Python", "The language engine"),
            ("Editor", "Where you write code"),
            ("Terminal", "Where you run and see results"),
        ],
        "script": "Reassure viewers: they do not need to master every tool. Today they only need Python, an editor, and a terminal.",
    },
    {
        "kind": "code",
        "section": "First code",
        "title": "First spell: make Python speak",
        "subtitle": "The smallest useful program is a message.",
        "code": "print(\"Hello, future Python builder!\")\nprint(\"Today I control the computer.\")",
        "callout": "print() sends text to the screen.",
        "script": "Run or demonstrate the code. Say the quotes hold text, and print sends that text to the screen.",
    },
    {
        "kind": "debug",
        "section": "Errors",
        "title": "Errors are not failure",
        "subtitle": "They are Python's way of pointing at the confusion.",
        "bad_code": "print(\"Hello Day 1\"",
        "good_code": "print(\"Hello Day 1\")",
        "script": "Show a missing bracket. Treat the error like a detective clue, not a personal insult.",
    },
    {
        "kind": "concept",
        "section": "Variables",
        "title": "Variables are name tags",
        "subtitle": "A variable gives a useful value a short name.",
        "points": [
            "channel = \"Code With Me\"",
            "episode = 1",
            "topic = \"Python basics\"",
        ],
        "script": "Use physical labels: a variable is like putting a sticker on a box so you can find the value later.",
    },
    {
        "kind": "code",
        "section": "Variables",
        "title": "Variables make code reusable",
        "subtitle": "Change one value, update the whole message.",
        "code": "channel = \"Code With Me\"\ntopic = \"Python\"\n\nprint(\"Welcome to\", channel)\nprint(\"Today we start\", topic)",
        "callout": "A name can be reused many times.",
        "script": "Point out that the same variable can appear in many places. This is the first taste of avoiding repetition.",
    },
    {
        "kind": "types",
        "section": "Data types",
        "title": "Python stores different kinds of things",
        "subtitle": "Do not memorise. Notice the shape.",
        "types": [
            ("str", "\"hello\"", "text"),
            ("int", "42", "whole number"),
            ("float", "3.14", "decimal"),
            ("bool", "True", "yes/no"),
        ],
        "script": "Explain types as shapes of values. For Day 1, recognition matters more than memorisation.",
    },
    {
        "kind": "code",
        "section": "Input",
        "title": "input() makes the program listen",
        "subtitle": "Now the viewer becomes part of the program.",
        "code": "name = input(\"What is your name? \")\nprint(\"Welcome,\", name)\nprint(\"Your Python journey starts now.\")",
        "callout": "input() pauses and waits for the user.",
        "script": "This is the emotional turning point. The program is no longer a poster; it is a conversation.",
    },
    {
        "kind": "lab",
        "section": "Mini project",
        "title": "Mini project: YouTube hook generator",
        "subtitle": "Build something your audience can instantly understand.",
        "code": "channel = input(\"Channel name: \")\ntopic = input(\"Video topic: \")\nvibe = input(\"Vibe: \")\n\nprint(\"Welcome to\", channel)\nprint(\"Today we make\", topic, \"feel\", vibe)\nprint(\"Watch till the end for your first Python win!\")",
        "script": "Frame this as creator-first Python. The first project should feel connected to their world, not like a random school exercise.",
    },
    {
        "kind": "activity",
        "section": "Practice",
        "title": "Pause-and-build challenge",
        "subtitle": "Give viewers a reason to touch the keyboard now.",
        "cards": [
            ("Level 1", "Print your name, city, and dream project."),
            ("Level 2", "Ask for a name and greet the viewer."),
            ("Level 3", "Ask for topic + vibe, then create a video intro line."),
        ],
        "script": "Ask viewers to pause the video and try one level. This creates participation, not passive watching.",
    },
    {
        "kind": "recap",
        "section": "Recap",
        "title": "Day 1 memory hooks",
        "subtitle": "If viewers remember only this, they can continue.",
        "points": [
            "print() = Python talks.",
            "variables = Python remembers.",
            "input() = Python listens.",
            "errors = Python gives clues.",
        ],
        "script": "Repeat the three-part model. This is the unique anchor of the episode: talks, remembers, listens.",
    },
    {
        "kind": "homework",
        "section": "Homework",
        "title": "Your Day 1 mission",
        "subtitle": "Build a tiny personal intro bot.",
        "tasks": [
            "Ask for name, city, and dream skill.",
            "Print a friendly intro in 3 lines.",
            "Comment your output under the video.",
        ],
        "script": "Convert homework into community engagement. Ask them to paste their output in comments.",
    },
    {
        "kind": "teaser",
        "section": "Next episode",
        "title": "Next: decisions",
        "subtitle": "We teach Python to choose with if / else.",
        "script": "End with a clear next step: Day 2 will make programs decide, like a mini brain.",
    },
]


def add_full_rect(slide, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def add_text(slide, text, left, top, width, height, size=28, color=None, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or PALETTE["ink"]
    return box


def add_pill(slide, text, left, top, width, color, text_color=None):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.42))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    frame = pill.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = text_color or PALETTE["navy"]
    return pill


def add_footer(slide, number, section):
    add_text(slide, f"{number:02d} / {len(SLIDES):02d}", 0.55, 7.02, 1.0, 0.25, 9, PALETTE["muted"])
    add_text(slide, section, 10.2, 7.02, 2.55, 0.25, 9, PALETTE["muted"], align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=None, dark=False):
    title_color = PALETTE["white"] if dark else PALETTE["ink"]
    subtitle_color = RGBColor(203, 213, 225) if dark else PALETTE["muted"]
    add_text(slide, title, 0.75, 0.72, 8.4, 0.75, 34, title_color, True)
    if subtitle:
        add_text(slide, subtitle, 0.78, 1.45, 9.0, 0.42, 15, subtitle_color)


def add_card(slide, title, body, left, top, width, height, accent):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = PALETTE["white"]
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.line.width = Pt(1)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.1), Inches(height))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    add_text(slide, title, left + 0.28, top + 0.18, width - 0.45, 0.35, 17, PALETTE["ink"], True)
    add_text(slide, body, left + 0.28, top + 0.64, width - 0.45, height - 0.75, 12.5, PALETTE["muted"])


def add_terminal(slide, code, left, top, width, height, title="demo.py"):
    terminal = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    terminal.fill.solid()
    terminal.fill.fore_color.rgb = PALETTE["code"]
    terminal.line.color.rgb = RGBColor(51, 65, 85)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
    bar.line.fill.background()
    for i, color in enumerate([PALETTE["red"], PALETTE["yellow"], PALETTE["green"]]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.18 + i * 0.25), Inches(top + 0.17), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
    add_text(slide, title, left + 0.78, top + 0.13, width - 1.0, 0.2, 9, RGBColor(203, 213, 225), font="Cascadia Mono")
    add_text(slide, code, left + 0.3, top + 0.72, width - 0.55, height - 0.9, 17, RGBColor(226, 232, 240), font="Cascadia Mono")


def add_gradient_dots(slide):
    specs = [
        (11.1, 0.55, 1.9, PALETTE["cyan"]),
        (10.6, 1.55, 1.1, PALETTE["violet"]),
        (12.1, 2.08, 0.6, PALETTE["pink"]),
        (0.3, 5.75, 0.55, PALETTE["yellow"]),
    ]
    for left, top, size, color in specs:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.fill.transparency = 18
        dot.line.fill.background()


def slide_cover(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["navy"])
    add_gradient_dots(slide)
    add_pill(slide, data["kicker"], 0.82, 0.68, 3.55, PALETTE["cyan"])
    add_text(slide, data["title"], 0.8, 1.55, 6.6, 0.8, 48, PALETTE["white"], True)
    add_text(slide, data["subtitle"], 0.83, 2.42, 6.9, 0.55, 26, PALETTE["cyan"], True)
    add_text(slide, "A Day 1 Python experience for viewers who want to build, not just watch.", 0.85, 3.35, 6.3, 0.9, 20, RGBColor(203, 213, 225))
    add_terminal(slide, ">>> print(\"I can make the computer listen\")\nI can make the computer listen", 7.5, 4.15, 4.8, 1.75, "first_moment.py")
    add_footer(slide, idx, data["section"])
    return slide


def slide_story(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["paper"])
    add_pill(slide, data["section"], 0.75, 0.38, 1.25, PALETTE["yellow"])
    add_title(slide, data["title"], data["subtitle"])
    colors = [PALETTE["blue"], PALETTE["violet"], PALETTE["pink"]]
    for i, (title, body) in enumerate(data["cards"]):
        add_card(slide, title, body, 0.85 + i * 4.05, 2.55, 3.45, 2.3, colors[i])
    add_text(slide, "Hook line for video: \"You are not learning a language first. You are learning how to give instructions.\"", 1.0, 5.45, 10.9, 0.55, 17, PALETTE["slate"], True, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, data["section"])
    return slide


def slide_map(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["navy"])
    add_gradient_dots(slide)
    add_pill(slide, data["section"], 0.75, 0.38, 1.2, PALETTE["cyan"])
    add_title(slide, data["title"], data["subtitle"], dark=True)
    for i, (num, title, body) in enumerate(data["steps"]):
        left = 1.1 + i * 4.0
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(2.55), Inches(1.0), Inches(1.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = [PALETTE["cyan"], PALETTE["yellow"], PALETTE["pink"]][i]
        circle.line.fill.background()
        add_text(slide, num, left, 2.76, 1.0, 0.25, 23, PALETTE["navy"], True, align=PP_ALIGN.CENTER)
        add_text(slide, title, left - 0.75, 3.82, 2.5, 0.35, 25, PALETTE["white"], True, align=PP_ALIGN.CENTER)
        add_text(slide, body, left - 0.9, 4.35, 2.8, 0.65, 13.5, RGBColor(203, 213, 225), align=PP_ALIGN.CENTER)
        if i < 2:
            add_text(slide, "->", left + 1.6, 2.86, 0.6, 0.3, 24, RGBColor(148, 163, 184), True, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, data["section"])
    return slide


def slide_concept(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["paper"])
    add_pill(slide, data["section"], 0.75, 0.38, 1.55, PALETTE["cyan"])
    add_title(slide, data["title"], data["subtitle"])
    for i, point in enumerate(data["points"]):
        top = 2.5 + i * 1.05
        icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(top), Inches(0.58), Inches(0.58))
        icon.fill.solid()
        icon.fill.fore_color.rgb = [PALETTE["blue"], PALETTE["violet"], PALETTE["green"]][i % 3]
        icon.line.fill.background()
        add_text(slide, str(i + 1), 1.05, top + 0.13, 0.58, 0.18, 13, PALETTE["white"], True, align=PP_ALIGN.CENTER)
        add_text(slide, point, 1.9, top + 0.05, 9.3, 0.45, 22, PALETTE["ink"], True)
    add_footer(slide, idx, data["section"])
    return slide


def slide_pipeline(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["paper"])
    add_pill(slide, data["section"], 0.75, 0.38, 1.35, PALETTE["yellow"])
    add_title(slide, data["title"], data["subtitle"])
    colors = [PALETTE["blue"], PALETTE["violet"], PALETTE["green"]]
    for i, (label, body) in enumerate(data["steps"]):
        left = 0.95 + i * 4.1
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.75), Inches(3.05), Inches(1.65))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()
        add_text(slide, label, left + 0.15, 3.0, 2.75, 0.35, 23, PALETTE["white"], True, align=PP_ALIGN.CENTER)
        add_text(slide, body, left + 0.25, 3.55, 2.55, 0.25, 13, PALETTE["white"], align=PP_ALIGN.CENTER)
        if i < 2:
            add_text(slide, "->", left + 3.25, 3.28, 0.55, 0.3, 26, PALETTE["muted"], True, align=PP_ALIGN.CENTER)
    add_terminal(slide, "$ python demo.py\nHello from Python!", 3.75, 5.15, 5.85, 1.25, "terminal")
    add_footer(slide, idx, data["section"])
    return slide


def slide_setup(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["navy"])
    add_gradient_dots(slide)
    add_pill(slide, data["section"], 0.75, 0.38, 1.15, PALETTE["green"])
    add_title(slide, data["title"], data["subtitle"], dark=True)
    for i, (title, body) in enumerate(data["items"]):
        left = 1.05 + i * 4.0
        add_card(slide, title, body, left, 2.75, 3.1, 2.0, [PALETTE["cyan"], PALETTE["yellow"], PALETTE["pink"]][i])
    add_text(slide, "Creator tip: show your screen from the first run, not only slides. Beginners trust what they can see.", 1.0, 5.85, 11.2, 0.42, 15, RGBColor(203, 213, 225), True, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, data["section"])
    return slide


def slide_code(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["paper"])
    add_pill(slide, data["section"], 0.75, 0.38, 1.55, PALETTE["cyan"])
    add_title(slide, data["title"], data["subtitle"])
    add_terminal(slide, data["code"], 0.95, 2.35, 7.25, 3.1)
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(2.55), Inches(3.45), Inches(2.25))
    callout.fill.solid()
    callout.fill.fore_color.rgb = PALETTE["navy"]
    callout.line.fill.background()
    add_text(slide, "Decode", 8.95, 2.88, 2.9, 0.35, 19, PALETTE["cyan"], True)
    add_text(slide, data["callout"], 8.95, 3.45, 2.75, 0.85, 17, PALETTE["white"], True)
    add_footer(slide, idx, data["section"])
    return slide


def slide_debug(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["paper"])
    add_pill(slide, data["section"], 0.75, 0.38, 1.15, PALETTE["red"], PALETTE["white"])
    add_title(slide, data["title"], data["subtitle"])
    add_terminal(slide, data["bad_code"], 0.95, 2.4, 5.45, 2.35, "broken.py")
    add_terminal(slide, data["good_code"], 6.95, 2.4, 5.45, 2.35, "fixed.py")
    add_text(slide, "Beginner mindset shift: an error is a map marker. It says, \"Look here first.\"", 1.1, 5.55, 11.1, 0.45, 18, PALETTE["slate"], True, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, data["section"])
    return slide


def slide_types(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["paper"])
    add_pill(slide, data["section"], 0.75, 0.38, 1.45, PALETTE["violet"], PALETTE["white"])
    add_title(slide, data["title"], data["subtitle"])
    colors = [PALETTE["blue"], PALETTE["green"], PALETTE["orange"], PALETTE["pink"]]
    for i, (name, value, meaning) in enumerate(data["types"]):
        left = 0.9 + (i % 2) * 6.1
        top = 2.45 + (i // 2) * 1.7
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.35), Inches(1.25))
        card.fill.solid()
        card.fill.fore_color.rgb = PALETTE["white"]
        card.line.color.rgb = RGBColor(226, 232, 240)
        add_pill(slide, name, left + 0.25, top + 0.2, 0.78, colors[i], PALETTE["white"])
        add_text(slide, value, left + 1.3, top + 0.25, 1.8, 0.3, 20, PALETTE["ink"], True, font="Cascadia Mono")
        add_text(slide, meaning, left + 3.05, top + 0.3, 1.8, 0.25, 13, PALETTE["muted"])
    add_footer(slide, idx, data["section"])
    return slide


def slide_lab(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["navy"])
    add_gradient_dots(slide)
    add_pill(slide, data["section"], 0.75, 0.38, 1.35, PALETTE["yellow"])
    add_title(slide, data["title"], data["subtitle"], dark=True)
    add_terminal(slide, data["code"], 0.95, 2.15, 7.45, 3.75, "creator_hook.py")
    add_text(slide, "Why this is unique", 8.85, 2.45, 3.2, 0.35, 20, PALETTE["cyan"], True)
    add_text(slide, "The first project is not a calculator. It connects Python to the viewer's own YouTube world.", 8.85, 3.05, 3.0, 1.35, 17, PALETTE["white"], True)
    add_footer(slide, idx, data["section"])
    return slide


def slide_activity(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["paper"])
    add_pill(slide, data["section"], 0.75, 0.38, 1.25, PALETTE["green"])
    add_title(slide, data["title"], data["subtitle"])
    colors = [PALETTE["green"], PALETTE["yellow"], PALETTE["pink"]]
    for i, (title, body) in enumerate(data["cards"]):
        add_card(slide, title, body, 0.85 + i * 4.05, 2.5, 3.45, 2.55, colors[i])
    add_text(slide, "On-camera prompt: \"Pause here. Choose one level. I will wait three seconds.\"", 1.0, 5.55, 11.2, 0.45, 18, PALETTE["slate"], True, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, data["section"])
    return slide


def slide_recap(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["navy"])
    add_gradient_dots(slide)
    add_pill(slide, data["section"], 0.75, 0.38, 1.15, PALETTE["cyan"])
    add_title(slide, data["title"], data["subtitle"], dark=True)
    for i, point in enumerate(data["points"]):
        add_card(slide, point.split(" = ")[0], point, 1.05 + (i % 2) * 5.75, 2.45 + (i // 2) * 1.55, 4.9, 1.05, [PALETTE["cyan"], PALETTE["yellow"], PALETTE["green"], PALETTE["pink"]][i])
    add_footer(slide, idx, data["section"])
    return slide


def slide_homework(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["paper"])
    add_pill(slide, data["section"], 0.75, 0.38, 1.35, PALETTE["yellow"])
    add_title(slide, data["title"], data["subtitle"])
    for i, task in enumerate(data["tasks"]):
        top = 2.35 + i * 1.0
        add_text(slide, f"{i + 1}.", 1.15, top, 0.4, 0.35, 22, PALETTE["blue"], True)
        add_text(slide, task, 1.75, top, 9.6, 0.35, 22, PALETTE["ink"], True)
    add_text(slide, "Comment prompt: \"Paste your bot output below. I will heart the cleanest ones.\"", 1.1, 5.65, 11.1, 0.45, 18, PALETTE["slate"], True, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, data["section"])
    return slide


def slide_teaser(prs, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, PALETTE["navy"])
    add_gradient_dots(slide)
    add_pill(slide, data["section"], 0.75, 0.38, 1.65, PALETTE["pink"], PALETTE["white"])
    add_text(slide, data["title"], 0.85, 1.8, 6.0, 0.8, 48, PALETTE["white"], True)
    add_text(slide, data["subtitle"], 0.9, 2.72, 6.7, 0.45, 22, PALETTE["cyan"], True)
    add_terminal(slide, "if viewer_is_ready:\n    print(\"See you in Day 2\")\nelse:\n    print(\"Replay Day 1 and build again\")", 7.1, 3.45, 4.95, 2.15, "day2_teaser.py")
    add_text(slide, "Like, subscribe, and bring your Day 1 output.", 0.92, 4.8, 5.7, 0.45, 20, RGBColor(203, 213, 225), True)
    add_footer(slide, idx, data["section"])
    return slide


BUILDERS = {
    "cover": slide_cover,
    "story": slide_story,
    "map": slide_map,
    "concept": slide_concept,
    "pipeline": slide_pipeline,
    "setup": slide_setup,
    "code": slide_code,
    "debug": slide_debug,
    "types": slide_types,
    "lab": slide_lab,
    "activity": slide_activity,
    "recap": slide_recap,
    "homework": slide_homework,
    "teaser": slide_teaser,
}


def write_notes():
    lines = [
        "# Python Intro / Day 1 - Speaker Notes",
        "",
        "Use these notes while recording the YouTube episode. The deck is designed around a different beginner experience: story, mental models, participation, and a creator-themed mini project.",
        "",
    ]
    for i, slide in enumerate(SLIDES, start=1):
        lines.extend(
            [
                f"## Slide {i}: {slide['title']}",
                f"- Section: {slide['section']}",
                f"- On-screen idea: {slide.get('subtitle', '')}",
                f"- Presenter note: {slide['script']}",
                "",
            ]
        )
    OUT_NOTES.write_text("\n".join(lines), encoding="utf-8")


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for idx, slide in enumerate(SLIDES, start=1):
        BUILDERS[slide["kind"]](prs, slide, idx)
    prs.save(OUT_PPTX)
    write_notes()


if __name__ == "__main__":
    build_deck()
    print(dedent(
        f"""
        Created:
        - {OUT_PPTX}
        - {OUT_NOTES}
        """
    ).strip())
