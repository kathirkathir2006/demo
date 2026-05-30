from pathlib import Path
from random import Random
from textwrap import wrap
from zipfile import ZIP_DEFLATED, ZipFile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path("Vallinate_Python_Day1_Industry_Edition.pptx")
OUT = Path("Vallinate_Python_Day1_Industry_Edition_PREMIUM.pptx")
PACKAGE = Path("Vallinate_Python_Day1_PREMIUM_Package.zip")
ASSET_DIR = Path("premium_assets")


COLORS = {
    "navy": RGBColor(6, 10, 28),
    "panel": RGBColor(12, 18, 42),
    "panel2": RGBColor(18, 26, 57),
    "white": RGBColor(248, 250, 252),
    "muted": RGBColor(148, 163, 184),
    "cyan": RGBColor(34, 211, 238),
    "blue": RGBColor(59, 130, 246),
    "violet": RGBColor(139, 92, 246),
    "pink": RGBColor(236, 72, 153),
    "green": RGBColor(34, 197, 94),
    "yellow": RGBColor(250, 204, 21),
    "orange": RGBColor(251, 146, 60),
    "red": RGBColor(248, 113, 113),
}


SLIDES = [
    {
        "tag": "INDUSTRY CONTEXT",
        "title": "Python in 2026: Why It Matters",
        "subtitle": "AI, cybersecurity, cloud, DevOps and automation all rely heavily on Python.",
        "kind": "network",
        "bullets": [
            "Python is a practical skill across modern tech teams.",
            "One language can open multiple career tracks.",
            "Day 1 starts with control: talk, remember, listen.",
        ],
    },
    {
        "tag": "MARKET REALITY",
        "title": "The Tech Market Reality",
        "subtitle": "Companies hire problem-solvers, not certificate collectors.",
        "kind": "dashboard",
        "bullets": [
            "Python appears in automation, AI, data and platform roles.",
            "Hiring teams value clear thinking and project evidence.",
            "Your learning goal: build small proof every day.",
        ],
    },
    {
        "tag": "WHERE IT IS USED",
        "title": "Where Python Is Used",
        "subtitle": "AI assistants, cloud automation, security tools, analytics and web apps.",
        "kind": "orbit",
        "bullets": [
            "AI: prompts, models, agents and data pipelines.",
            "Cloud and DevOps: scripts, APIs and infrastructure helpers.",
            "Cybersecurity: scanning, parsing logs and building tools.",
        ],
    },
    {
        "tag": "CAREER MAP",
        "title": "Career Paths",
        "subtitle": "Python can become your common language across tech specializations.",
        "kind": "careers",
        "bullets": [
            "AI Engineer",
            "Cybersecurity Analyst",
            "DevOps Engineer",
            "Cloud Engineer",
            "Data Engineer",
        ],
    },
    {
        "tag": "DAY 1 FRAMEWORK",
        "title": "What You Will Learn",
        "subtitle": "Talk. Remember. Listen.",
        "kind": "three_steps",
        "bullets": [
            "Talk: print() shows output.",
            "Remember: variables store values.",
            "Listen: input() collects user data.",
        ],
    },
    {
        "tag": "HOW IT RUNS",
        "title": "How Python Works",
        "subtitle": "You write code -> Python interprets it -> the computer executes it.",
        "kind": "pipeline",
        "bullets": [
            "The interpreter reads your code.",
            "Instructions run in order.",
            "The result appears in the terminal or app.",
        ],
    },
    {
        "tag": "FIRST PROGRAM",
        "title": "First Program",
        "subtitle": "The smallest useful program is a message on the screen.",
        "kind": "code",
        "code": "print('Hello Vallinate')",
        "bullets": [
            "print() is your first output command.",
            "Quotes hold text.",
            "Parentheses carry the message into the function.",
        ],
    },
    {
        "tag": "DEBUG MINDSET",
        "title": "Errors Are Feedback",
        "subtitle": "Errors help developers locate problems quickly.",
        "kind": "error",
        "code": "print('Hello Vallinate'\n# missing closing parenthesis",
        "bullets": [
            "An error is not a failure.",
            "Read the line number first.",
            "Fix one problem, then run again.",
        ],
    },
    {
        "tag": "VARIABLES",
        "title": "Variables",
        "subtitle": "Applications remember users, products and settings using stored data.",
        "kind": "database",
        "code": "name = 'Vallinate'\nrole = 'AI Engineer'\nprint(name, role)",
        "bullets": [
            "A variable is a name tag for a value.",
            "Use readable names.",
            "Change the value once, reuse it many times.",
        ],
    },
    {
        "tag": "DATA TYPES",
        "title": "Data Types",
        "subtitle": "String, integer, float and boolean with real-world examples.",
        "kind": "types",
        "bullets": [
            "str: text like a name or city.",
            "int: whole numbers like age or score.",
            "float: decimals like price or rating.",
            "bool: True or False decisions.",
        ],
    },
    {
        "tag": "INPUT",
        "title": "Input",
        "subtitle": "Collect information from users like login and signup forms.",
        "kind": "form",
        "code": "name = input('Your name: ')\nprint('Welcome', name)",
        "bullets": [
            "input() pauses the program.",
            "The user types an answer.",
            "Python stores that answer as text.",
        ],
    },
    {
        "tag": "MINI PROJECT",
        "title": "Mini Project",
        "subtitle": "Career Path Generator: ask name and dream role, then display a personalized message.",
        "kind": "project",
        "code": "name = input('Name: ')\nrole = input('Dream role: ')\nprint(name, 'can start the path to', role)",
        "bullets": [
            "This feels like a real app interaction.",
            "It combines print(), variables and input().",
            "It gives learners a Day 1 win.",
        ],
    },
    {
        "tag": "INDUSTRY CHALLENGE",
        "title": "Industry Challenge",
        "subtitle": "Modify the project for AI, cybersecurity, cloud or DevOps careers.",
        "kind": "challenge",
        "bullets": [
            "Ask for one target career.",
            "Print a custom first project idea.",
            "Share the output in the comments.",
        ],
    },
    {
        "tag": "90-DAY ROADMAP",
        "title": "90-Day Roadmap",
        "subtitle": "Basics -> projects -> automation -> specialization.",
        "kind": "roadmap",
        "bullets": [
            "Basics: syntax, variables, conditions and loops.",
            "Projects: small apps with visible outcomes.",
            "Automation: files, APIs and repeatable workflows.",
            "Specialization: AI, cyber, cloud, DevOps or data.",
        ],
    },
    {
        "tag": "NEXT EPISODE",
        "title": "Next Episode",
        "subtitle": "if, else and decision making.",
        "kind": "next",
        "code": "if ready:\n    print('Day 2 begins')\nelse:\n    print('Rewatch and build again')",
        "bullets": [
            "Next, Python will choose between paths.",
            "Decision making turns scripts into smarter programs.",
            "Bring your mini project output.",
        ],
    },
]


def font(size, bold=False):
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def hex_to_rgb(value):
    value = value.lstrip("#")
    return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))


def gradient(size, start="#050816", end="#172554"):
    w, h = size
    img = Image.new("RGB", size, start)
    draw = ImageDraw.Draw(img)
    sr, sg, sb = hex_to_rgb(start)
    er, eg, eb = hex_to_rgb(end)
    for y in range(h):
        ratio = y / max(h - 1, 1)
        r = int(sr + (er - sr) * ratio)
        g = int(sg + (eg - sg) * ratio)
        b = int(sb + (eb - sb) * ratio)
        draw.line([(0, y), (w, y)], fill=(r, g, b))
    return img


def rounded(draw, xy, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def draw_text(draw, text, xy, size=32, fill=(248, 250, 252), bold=False, max_width=None, line_gap=8):
    fnt = font(size, bold)
    if not max_width:
        draw.text(xy, text, font=fnt, fill=fill)
        return
    chars = max(12, int(max_width / (size * 0.55)))
    lines = []
    for piece in text.split("\n"):
        lines.extend(wrap(piece, width=chars) or [""])
    x, y = xy
    for line in lines:
        draw.text((x, y), line, font=fnt, fill=fill)
        y += size + line_gap


def add_glow(draw, center, color, max_radius):
    x, y = center
    for radius in range(max_radius, 8, -18):
        alpha = radius / max_radius
        fill = tuple(int(c * alpha + 10 * (1 - alpha)) for c in color)
        draw.ellipse((x - radius, y - radius, x + radius, y + radius), outline=fill, width=3)


def terminal(draw, xy, size, lines):
    x, y = xy
    w, h = size
    rounded(draw, (x, y, x + w, y + h), 28, (8, 13, 31), (51, 65, 85), 3)
    rounded(draw, (x, y, x + w, y + 70), 28, (15, 23, 42))
    for i, c in enumerate([(248, 113, 113), (250, 204, 21), (34, 197, 94)]):
        draw.ellipse((x + 30 + i * 34, y + 27, x + 46 + i * 34, y + 43), fill=c)
    draw_text(draw, "demo.py", (x + 145, y + 21), 22, (203, 213, 225), False)
    yy = y + 105
    for line in lines:
        draw_text(draw, line, (x + 38, yy), 30, (226, 232, 240), False)
        yy += 44


def base_asset(slide_no, title):
    img = gradient((1200, 720), "#050816", "#111c44")
    draw = ImageDraw.Draw(img)
    rnd = Random(slide_no)
    for _ in range(60):
        x, y = rnd.randint(0, 1200), rnd.randint(0, 720)
        color = rnd.choice([(34, 211, 238), (139, 92, 246), (236, 72, 153), (59, 130, 246)])
        draw.ellipse((x, y, x + 3, y + 3), fill=color)
    add_glow(draw, (930, 160), (34, 211, 238), 230)
    add_glow(draw, (240, 620), (139, 92, 246), 210)
    draw_text(draw, f"{slide_no:02d}", (50, 45), 42, (34, 211, 238), True)
    draw_text(draw, title, (50, 98), 34, (248, 250, 252), True, max_width=480)
    return img, draw


def save_asset(slide_no, data):
    img, draw = base_asset(slide_no, data["title"])
    kind = data["kind"]
    if kind == "network":
        centers = [(720, 210), (910, 150), (1010, 330), (790, 420), (590, 350)]
        for a in centers:
            for b in centers:
                if a != b:
                    draw.line((a, b), fill=(51, 65, 85), width=2)
        labels = ["AI", "CYBER", "CLOUD", "DEVOPS", "DATA"]
        for (x, y), label in zip(centers, labels):
            draw.ellipse((x - 62, y - 62, x + 62, y + 62), fill=(12, 18, 42), outline=(34, 211, 238), width=4)
            draw_text(draw, label, (x - 42, y - 15), 25, (248, 250, 252), True)
        terminal(draw, (90, 430), (520, 210), [">>> python", "builds careers"])
    elif kind == "dashboard":
        for i, label in enumerate(["AI", "Cloud", "Cyber", "Data"]):
            x = 590 + i * 135
            h = [250, 190, 220, 160][i]
            rounded(draw, (x, 520 - h, x + 85, 520), 18, [(34, 211, 238), (139, 92, 246), (236, 72, 153), (34, 197, 94)][i])
            draw_text(draw, label, (x + 5, 545), 23, (203, 213, 225), True)
        rounded(draw, (110, 420, 500, 610), 30, (12, 18, 42), (34, 211, 238), 3)
        draw_text(draw, "Project proof > passive learning", (140, 480), 34, (248, 250, 252), True, 320)
    elif kind == "orbit":
        cx, cy = 850, 350
        draw.ellipse((cx - 95, cy - 95, cx + 95, cy + 95), fill=(34, 211, 238), outline=(248, 250, 252), width=4)
        draw_text(draw, "PY", (cx - 39, cy - 35), 54, (5, 8, 22), True)
        labels = ["AI", "Cloud", "Security", "Web", "Data"]
        for i, label in enumerate(labels):
            angle = i * 72
            x = cx + int(280 * __import__("math").cos(__import__("math").radians(angle)))
            y = cy + int(190 * __import__("math").sin(__import__("math").radians(angle)))
            draw.line((cx, cy, x, y), fill=(71, 85, 105), width=3)
            rounded(draw, (x - 82, y - 32, x + 82, y + 32), 22, (15, 23, 42), (139, 92, 246), 3)
            draw_text(draw, label, (x - 55, y - 15), 22, (248, 250, 252), True)
    elif kind == "careers":
        names = ["AI", "Cyber", "DevOps", "Cloud", "Data"]
        colors = [(34, 211, 238), (236, 72, 153), (250, 204, 21), (59, 130, 246), (34, 197, 94)]
        for i, name in enumerate(names):
            x = 530 + (i % 3) * 210
            y = 220 + (i // 3) * 180
            rounded(draw, (x, y, x + 175, y + 110), 30, colors[i])
            draw_text(draw, name, (x + 28, y + 36), 30, (5, 8, 22), True)
        draw.line((620, 520, 980, 520), fill=(34, 211, 238), width=5)
        draw_text(draw, "One foundation. Many paths.", (600, 565), 32, (248, 250, 252), True)
    elif kind == "three_steps":
        steps = [("Talk", "print()"), ("Remember", "variables"), ("Listen", "input()")]
        for i, (a, b) in enumerate(steps):
            x = 530 + i * 210
            rounded(draw, (x, 250, x + 170, 430), 36, [(34, 211, 238), (250, 204, 21), (236, 72, 153)][i])
            draw_text(draw, a, (x + 24, 295), 31, (5, 8, 22), True)
            draw_text(draw, b, (x + 24, 350), 24, (5, 8, 22), True)
    elif kind == "pipeline":
        nodes = [("Write", 520), ("Interpret", 760), ("Execute", 1000)]
        for label, x in nodes:
            draw.ellipse((x - 75, 285, x + 75, 435), fill=(15, 23, 42), outline=(34, 211, 238), width=5)
            draw_text(draw, label, (x - 55, 348), 24, (248, 250, 252), True)
        draw.line((595, 360, 685, 360), fill=(250, 204, 21), width=6)
        draw.line((835, 360, 925, 360), fill=(250, 204, 21), width=6)
    elif kind in {"code", "error", "database", "form", "project", "next"}:
        code_lines = data.get("code", "print('Hello')")
        terminal(draw, (520, 230), (590, 310), code_lines.splitlines())
        if kind == "error":
            draw.line((950, 175, 1080, 305), fill=(248, 113, 113), width=14)
            draw.line((1080, 175, 950, 305), fill=(248, 113, 113), width=14)
        elif kind == "database":
            for y in [180, 235, 290]:
                draw.ellipse((880, y, 1080, y + 70), fill=(18, 26, 57), outline=(34, 197, 94), width=4)
        elif kind == "form":
            rounded(draw, (620, 140, 1050, 540), 30, (248, 250, 252))
            for y in [230, 320, 410]:
                rounded(draw, (670, y, 1000, y + 48), 16, (226, 232, 240))
            draw_text(draw, "USER INPUT", (700, 165), 30, (15, 23, 42), True)
        elif kind == "project":
            draw.polygon([(880, 140), (1060, 260), (995, 530), (760, 520), (700, 260)], fill=(34, 211, 238), outline=(248, 250, 252))
            draw_text(draw, "CAREER\nGENERATOR", (790, 285), 34, (5, 8, 22), True)
        elif kind == "next":
            draw.polygon([(780, 170), (1020, 350), (780, 530)], fill=(236, 72, 153))
    elif kind == "types":
        items = [("str", "text"), ("int", "42"), ("float", "3.14"), ("bool", "True")]
        for i, (a, b) in enumerate(items):
            x = 530 + (i % 2) * 280
            y = 230 + (i // 2) * 155
            rounded(draw, (x, y, x + 230, y + 105), 26, [(34, 211, 238), (250, 204, 21), (139, 92, 246), (34, 197, 94)][i])
            draw_text(draw, a, (x + 25, y + 18), 35, (5, 8, 22), True)
            draw_text(draw, b, (x + 25, y + 65), 25, (5, 8, 22), True)
    elif kind == "challenge":
        for i, label in enumerate(["AI", "Cyber", "Cloud", "DevOps"]):
            x = 540 + i * 145
            draw.ellipse((x, 250, x + 110, 360), fill=[(34, 211, 238), (236, 72, 153), (59, 130, 246), (250, 204, 21)][i])
            draw_text(draw, label, (x + 15, 395), 24, (248, 250, 252), True)
        draw_text(draw, "Choose one. Customize the project.", (560, 520), 34, (248, 250, 252), True)
    elif kind == "roadmap":
        labels = ["Basics", "Projects", "Automation", "Specialize"]
        for i, label in enumerate(labels):
            x = 470 + i * 175
            y = 245 + (i % 2) * 115
            draw.line((x + 80, y + 50, x + 215, 300 + ((i + 1) % 2) * 115), fill=(34, 211, 238), width=5)
            rounded(draw, (x, y, x + 160, y + 95), 28, [(34, 211, 238), (139, 92, 246), (34, 197, 94), (236, 72, 153)][i])
            draw_text(draw, label, (x + 15, y + 32), 23, (5, 8, 22), True)
    path = ASSET_DIR / f"slide_{slide_no:02d}_{kind}.png"
    img.save(path)
    return path


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["navy"]
    bg.line.fill.background()


def add_text(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, font_name="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_pill(slide, text, x, y, w, color):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    frame = pill.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = COLORS["navy"]


def add_panel(slide, x, y, w, h):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["panel"]
    panel.line.color.rgb = RGBColor(30, 41, 76)
    panel.line.width = Pt(1.1)
    return panel


def add_bullets(slide, bullets, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(9)
        p.text = f"- {bullet}"
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(203, 213, 225)


def add_code_card(slide, code, x, y, w, h):
    add_panel(slide, x, y, w, h)
    add_text(slide, "demo.py", x + 0.25, y + 0.2, w - 0.5, 0.25, 10, COLORS["cyan"], True, font_name="Cascadia Mono")
    add_text(slide, code, x + 0.25, y + 0.62, w - 0.5, h - 0.75, 15, COLORS["white"], False, font_name="Cascadia Mono")


def build_deck(asset_paths):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    accents = [COLORS["cyan"], COLORS["yellow"], COLORS["violet"], COLORS["green"], COLORS["pink"]]
    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        slide.shapes.add_picture(str(asset_paths[idx - 1]), Inches(6.65), Inches(0), Inches(6.68), Inches(7.5))
        add_pill(slide, data["tag"], 0.65, 0.5, min(2.5, 0.12 * len(data["tag"]) + 0.45), accents[(idx - 1) % len(accents)])
        add_text(slide, data["title"], 0.65, 1.1, 5.55, 1.05, 31, COLORS["white"], True)
        add_text(slide, data["subtitle"], 0.68, 2.12, 5.45, 0.75, 16, COLORS["muted"])
        add_panel(slide, 0.68, 3.15, 5.35, 2.05)
        add_bullets(slide, data["bullets"], 0.95, 3.42, 4.8, 1.55)
        if "code" in data:
            add_code_card(slide, data["code"], 0.68, 5.45, 5.35, 1.15)
        else:
            add_text(slide, "Premium teaching hook", 0.82, 5.58, 4.8, 0.25, 12, COLORS["cyan"], True)
            add_text(slide, "Connect every concept to real-world work so beginners feel the purpose immediately.", 0.82, 5.92, 4.85, 0.55, 13, COLORS["white"], True)
        add_text(slide, f"{idx:02d} / {len(SLIDES):02d}", 0.68, 7.05, 1.0, 0.2, 9, COLORS["muted"])
        add_text(slide, "Vallinate Python Day 1", 4.35, 7.05, 1.7, 0.2, 9, COLORS["muted"], align=PP_ALIGN.RIGHT)
    prs.save(OUT)


def package_files():
    with ZipFile(PACKAGE, "w", ZIP_DEFLATED) as zf:
        zf.write(OUT, OUT.name)
        zf.write(SOURCE, SOURCE.name)


def main():
    if not SOURCE.exists():
        raise SystemExit(f"Source PPT not found: {SOURCE}")
    ASSET_DIR.mkdir(exist_ok=True)
    asset_paths = [save_asset(i, slide) for i, slide in enumerate(SLIDES, start=1)]
    build_deck(asset_paths)
    package_files()
    print(f"Created {OUT}")
    print(f"Created {PACKAGE}")
    print(f"Created {len(asset_paths)} image assets in {ASSET_DIR}")


if __name__ == "__main__":
    main()
